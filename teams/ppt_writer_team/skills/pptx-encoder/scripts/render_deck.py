#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import math
import sys
from pathlib import Path
from typing import Any


def fail_missing_dependency() -> None:
    print(
        "Missing dependency: python-pptx. Install with `pip install python-pptx`.",
        file=sys.stderr,
    )
    raise SystemExit(2)


try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    fail_missing_dependency()


DEFAULT_STYLE: dict[str, Any] = {
    "slide_size": "16:9",
    "theme_colors": {
        "background": "FFFFFF",
        "text": "1F2937",
        "primary": "2563EB",
        "secondary": "14B8A6",
        "accent": "F59E0B",
        "muted": "6B7280",
        "surface": "F3F4F6",
    },
    "fonts": {"title": "Arial", "body": "Arial", "caption": "Arial"},
    "typography": {
        "title_size": 30,
        "subtitle_size": 18,
        "body_size": 15,
        "caption_size": 9,
        "min_title_size": 20,
        "min_body_size": 10,
    },
    "layout_tokens": {
        "margin_left": 0.65,
        "margin_right": 0.65,
        "header_height": 0.25,
        "footer_height": 0.28,
        "gutter": 0.28,
    },
    "template_patterns": {
        "uses_header": False,
        "uses_footer": True,
        "header_text": "",
        "footer_text": "",
    },
}


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    if not isinstance(data, dict):
        raise ValueError(f"{path} must contain a JSON object")
    return data


def merge_style(style: dict[str, Any]) -> dict[str, Any]:
    merged = json.loads(json.dumps(DEFAULT_STYLE))
    for key, value in style.items():
        if isinstance(value, dict) and isinstance(merged.get(key), dict):
            merged[key].update(value)
        else:
            merged[key] = value
    return merged


def template_fidelity(style: dict[str, Any]) -> dict[str, Any]:
    value = style.get("template_fidelity", {})
    return value if isinstance(value, dict) else {}


def use_template_chrome(style: dict[str, Any]) -> bool:
    fidelity = template_fidelity(style)
    return bool(fidelity.get("preserve_master")) and str(fidelity.get("mode", "")).lower() in {"strict", "adaptive"}


def color(style: dict[str, Any], name: str) -> RGBColor:
    value = style.get("theme_colors", {}).get(name, name)
    if not isinstance(value, str):
        value = "000000"
    value = value.strip().lstrip("#")
    if len(value) != 6:
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def inches(value: float):
    return Inches(value)


def as_inches(value) -> float:
    return float(value) / 914400


def text_length(text: str | list[Any]) -> int:
    if isinstance(text, list):
        return sum(len(str(item)) for item in text) + max(0, len(text) - 1) * 3
    return len(str(text))


def fitted_font_size(
    text: str | list[Any],
    *,
    width,
    height,
    preferred: int,
    minimum: int,
    maximum: int | None = None,
) -> int:
    maximum = maximum or preferred
    width_in = max(as_inches(width), 0.2)
    height_in = max(as_inches(height), 0.2)
    raw_text = "\n".join(str(item) for item in text) if isinstance(text, list) else str(text)
    line_count = max(1, raw_text.count("\n") + 1)
    chars = max(1, text_length(text))

    size = min(preferred, maximum)
    while size > minimum:
        chars_per_line = max(8, int(width_in * 144 / max(size, 1)))
        estimated_lines = max(line_count, math.ceil(chars / chars_per_line))
        estimated_height = estimated_lines * size * 1.22 / 72
        if estimated_height <= height_in * 0.9:
            break
        size -= 1
    return max(minimum, size)


def estimated_line_count(text: str | list[Any], *, width, font_size: int) -> int:
    width_in = max(as_inches(width), 0.2)
    raw_text = "\n".join(str(item) for item in text) if isinstance(text, list) else str(text)
    line_count = max(1, raw_text.count("\n") + 1)
    chars_per_line = max(8, int(width_in * 144 / max(font_size, 1)))
    return max(line_count, math.ceil(max(1, text_length(text)) / chars_per_line))


def warn_if_dense(
    warnings: list[str],
    *,
    slide_index: int,
    label: str,
    text: str | list[Any],
    width,
    height,
    font_size: int,
    minimum: int,
) -> None:
    if font_size > minimum:
        return
    lines = estimated_line_count(text, width=width, font_size=font_size)
    needed_height = lines * font_size * 1.22 / 72
    if needed_height > as_inches(height) * 0.9:
        warnings.append(f"Slide {slide_index}: `{label}` is too dense at minimum font size; split or shorten this content.")


def set_run(run, *, font_name: str, size: int, rgb: RGBColor, bold: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = rgb
    run.font.bold = bold


def text_color_for_role(style: dict[str, Any], role: str) -> RGBColor:
    if role in {"footer", "caption"}:
        return color(style, "muted")
    return color(style, "text")


def add_textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    style,
    *,
    role: str = "body",
    size: int | None = None,
    bold: bool | None = None,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    fonts = style.get("fonts", {})
    typography = style.get("typography", {})
    font_key = "title" if role == "title" else "caption" if role in {"caption", "footer"} else "body"
    font_name = fonts.get(font_key, fonts.get("body", "Arial"))
    preferred = size or int(typography.get(f"{role}_size", typography.get("body_size", 15)))
    minimum = int(typography.get("min_title_size" if role == "title" else "min_body_size", 10))
    if role in {"caption", "footer"}:
        minimum = min(minimum, 8)
    font_size = fitted_font_size(text, width=width, height=height, preferred=preferred, minimum=minimum)
    set_run(
        run,
        font_name=font_name,
        size=font_size,
        rgb=text_color_for_role(style, role),
        bold=(role == "title") if bold is None else bold,
    )
    return box


def add_hyperlink_textbox(
    slide,
    text: str,
    url: str,
    left,
    top,
    width,
    height,
    style,
    *,
    role: str = "body",
    size: int | None = None,
    bold: bool | None = None,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    if url:
        run.hyperlink.address = url
    fonts = style.get("fonts", {})
    typography = style.get("typography", {})
    font_key = "title" if role == "title" else "caption" if role in {"caption", "footer"} else "body"
    font_name = fonts.get(font_key, fonts.get("body", "Arial"))
    preferred = size or int(typography.get(f"{role}_size", typography.get("body_size", 15)))
    minimum = int(typography.get("min_title_size" if role == "title" else "min_body_size", 10))
    font_size = fitted_font_size(text, width=width, height=height, preferred=preferred, minimum=minimum)
    set_run(
        run,
        font_name=font_name,
        size=font_size,
        rgb=color(style, "primary"),
        bold=(role == "title") if bold is None else bold,
    )
    return box


def add_bullets(slide, items: list[str], left, top, width, height, style, *, size: int | None = None) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.clear()
    fonts = style.get("fonts", {})
    typography = style.get("typography", {})
    preferred = size or int(typography.get("body_size", 15))
    font_size = fitted_font_size(items, width=width, height=height, preferred=preferred, minimum=int(typography.get("min_body_size", 10)))
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.level = 0
        paragraph.space_after = Pt(3)
        run = paragraph.add_run()
        run.text = f"- {str(item)}"
        set_run(
            run,
            font_name=fonts.get("body", "Arial"),
            size=font_size,
            rgb=color(style, "text"),
        )


def add_table(slide, rows: list[list[Any]], left, top, width, height, style) -> None:
    if not rows:
        return
    row_count = len(rows)
    col_count = max(len(row) for row in rows)
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = table_shape.table
    font_size = fitted_font_size(
        [" ".join(str(cell) for cell in row) for row in rows],
        width=width,
        height=height,
        preferred=12 if row_count <= 6 else 10,
        minimum=8,
    )
    for row_index, row in enumerate(rows):
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            cell.text = str(row[col_index]) if col_index < len(row) else ""
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = color(style, "surface" if row_index == 0 else "background")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_run(
                        run,
                        font_name=style.get("fonts", {}).get("body", "Arial"),
                        size=font_size,
                        rgb=color(style, "text"),
                        bold=row_index == 0,
                    )


def add_placeholder(slide, label: str, left, top, width, height, style, *, kind: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(style, "surface")
    shape.line.color.rgb = color(style, "muted")
    add_textbox(slide, f"{kind}: {label}", left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3), style, size=13)


def resolve_local_path(path_value: str, *, base_dir: Path) -> Path | None:
    if not path_value:
        return None
    candidate = Path(path_value).expanduser()
    candidates = [candidate]
    if not candidate.is_absolute():
        candidates.append(base_dir / candidate)
        candidates.append(base_dir.parent / candidate)
        candidates.append(Path.cwd() / candidate)
    for path in candidates:
        if path.exists() and path.is_file():
            return path
    return None


def resolve_style_path(path_value: str, *, base_dir: Path) -> Path | None:
    if not path_value:
        return None
    candidate = Path(path_value).expanduser()
    candidates = [candidate]
    if not candidate.is_absolute():
        candidates.append(base_dir / candidate)
        candidates.append(base_dir.parent / candidate)
        candidates.append(Path.cwd() / candidate)
    for path in candidates:
        if path.exists() and path.is_file():
            return path
    return None


def image_path_from_element(element: dict[str, Any]) -> str:
    for key in ("processed_path", "image_path", "path", "src", "source"):
        value = element.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    return ""


def add_image(slide, image_path: Path, left, top, width, height, style, *, caption: str = "") -> None:
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    if caption:
        caption_height = Inches(0.22)
        add_textbox(
            slide,
            caption,
            left,
            top + height - caption_height,
            width,
            caption_height,
            style,
            role="caption",
            size=8,
        )


def set_click_url(shape, url: str) -> None:
    if not url:
        return
    try:
        shape.click_action.hyperlink.address = url
    except Exception:
        pass


def add_video_link(
    slide,
    element: dict[str, Any],
    left,
    top,
    width,
    height,
    style,
    *,
    base_dir: Path,
) -> None:
    title = str(element.get("title") or element.get("label") or "Video")
    url = str(element.get("video_url") or element.get("url") or element.get("href") or "")
    thumbnail_value = str(element.get("thumbnail_path") or element.get("image_path") or element.get("processed_path") or "")
    thumbnail = resolve_local_path(thumbnail_value, base_dir=base_dir)

    if thumbnail:
        picture = slide.shapes.add_picture(str(thumbnail), left, top, width=width, height=height)
        set_click_url(picture, url)
        overlay_h = Inches(0.42)
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height - overlay_h, width, overlay_h)
        banner.fill.solid()
        banner.fill.fore_color.rgb = color(style, "surface")
        banner.line.color.rgb = color(style, "primary")
        set_click_url(banner, url)
        add_hyperlink_textbox(
            slide,
            f"Play video: {title}",
            url,
            left + Inches(0.12),
            top + height - overlay_h + Inches(0.06),
            width - Inches(0.24),
            overlay_h - Inches(0.08),
            style,
            size=10,
            bold=True,
        )
    else:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = color(style, "surface")
        card.line.color.rgb = color(style, "primary")
        set_click_url(card, url)
        add_textbox(slide, "VIDEO", left + Inches(0.18), top + Inches(0.18), width - Inches(0.36), Inches(0.28), style, role="caption", size=9, bold=True)
        add_hyperlink_textbox(slide, title, url, left + Inches(0.18), top + Inches(0.58), width - Inches(0.36), height - Inches(1.0), style, size=15, bold=True)
        if url:
            add_hyperlink_textbox(slide, url, url, left + Inches(0.18), top + height - Inches(0.38), width - Inches(0.36), Inches(0.24), style, role="caption", size=8)

    caption = str(element.get("caption") or "")
    if caption:
        add_textbox(slide, caption, left, top + height - Inches(0.22), width, Inches(0.22), style, role="caption", size=8)


def add_notes(slide, notes: str) -> None:
    if not notes:
        return
    try:
        text_frame = slide.notes_slide.notes_text_frame
        text_frame.text = notes
    except Exception:
        pass


def add_design_chrome(slide, style: dict[str, Any], index: int, total: int, headline: str) -> None:
    if use_template_chrome(style):
        return
    tokens = style.get("layout_tokens", {})
    patterns = style.get("template_patterns", {})
    margin_left = float(tokens.get("margin_left", 0.65))
    margin_right = float(tokens.get("margin_right", 0.65))

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = color(style, "primary")
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.36), Inches(13.333), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color(style, "secondary")
    accent.line.fill.background()

    if patterns.get("uses_header") or patterns.get("header_text"):
        header = str(patterns.get("header_text") or headline)
        add_textbox(
            slide,
            header[:90],
            inches(margin_left),
            inches(0.12),
            inches(13.333 - margin_left - margin_right - 0.8),
            inches(0.22),
            style,
            role="caption",
            size=8,
        )

    footer_text = str(patterns.get("footer_text") or "")
    if footer_text:
        add_textbox(slide, footer_text[:100], inches(margin_left), inches(6.98), inches(7.0), inches(0.22), style, role="footer", size=8)

    add_textbox(
        slide,
        f"{index}/{total}",
        inches(12.1),
        inches(6.98),
        inches(0.58),
        inches(0.22),
        style,
        role="footer",
        size=8,
        align=PP_ALIGN.RIGHT,
    )


def content_boxes(element_count: int, *, top: float, bottom: float, left: float, right: float, gutter: float) -> list[tuple[Any, Any, Any, Any]]:
    available_w = right - left
    available_h = bottom - top
    if element_count <= 1:
        return [(inches(left), inches(top), inches(available_w), inches(available_h))]
    if element_count == 2:
        col_w = (available_w - gutter) / 2
        return [
            (inches(left), inches(top), inches(col_w), inches(available_h)),
            (inches(left + col_w + gutter), inches(top), inches(col_w), inches(available_h)),
        ]

    rows = math.ceil(element_count / 2)
    col_w = (available_w - gutter) / 2
    row_h = (available_h - gutter * (rows - 1)) / rows
    boxes = []
    for idx in range(element_count):
        col = idx % 2
        row = idx // 2
        boxes.append((inches(left + col * (col_w + gutter)), inches(top + row * (row_h + gutter)), inches(col_w), inches(row_h)))
    return boxes


def slot_to_box(slot: dict[str, Any]) -> tuple[Any, Any, Any, Any] | None:
    try:
        return (
            inches(float(slot["x"])),
            inches(float(slot["y"])),
            inches(float(slot["w"])),
            inches(float(slot["h"])),
        )
    except Exception:
        return None


def layout_slots(style: dict[str, Any], layout: str) -> list[dict[str, Any]]:
    layouts = style.get("layouts", {})
    if not isinstance(layouts, dict):
        return []
    layout_spec = layouts.get(layout, {})
    if not isinstance(layout_spec, dict):
        return []
    slots = layout_spec.get("slots", [])
    return [slot for slot in slots if isinstance(slot, dict)]


def box_for_role(style: dict[str, Any], layout: str, role: str) -> tuple[Any, Any, Any, Any] | None:
    for slot in layout_slots(style, layout):
        if str(slot.get("role", "")).lower() == role:
            return slot_to_box(slot)
    return None


def element_boxes_from_template(style: dict[str, Any], layout: str, elements: list[dict[str, Any]]) -> list[tuple[Any, Any, Any, Any]] | None:
    slots = [slot for slot in layout_slots(style, layout) if str(slot.get("role", "")).lower() not in {"headline", "title", "footer", "header"}]
    if not slots:
        return None
    boxes: list[tuple[Any, Any, Any, Any]] = []
    used: set[int] = set()
    for element in elements:
        requested = str(element.get("slot") or element.get("role") or element.get("type") or "").lower()
        match_index = None
        for idx, slot in enumerate(slots):
            if idx in used:
                continue
            slot_role = str(slot.get("role", "")).lower()
            if requested and requested == slot_role:
                match_index = idx
                break
        if match_index is None:
            for idx, _slot in enumerate(slots):
                if idx not in used:
                    match_index = idx
                    break
        if match_index is None:
            return None
        box = slot_to_box(slots[match_index])
        if box is None:
            return None
        used.add(match_index)
        boxes.append(box)
    return boxes


def render_element(
    slide,
    element: dict[str, Any],
    box: tuple[Any, Any, Any, Any],
    style: dict[str, Any],
    warnings: list[str],
    index: int,
    *,
    base_dir: Path,
) -> None:
    left, top, width, height = box
    element_type = element.get("type", "text")
    typography = style.get("typography", {})
    min_body = int(typography.get("min_body_size", 10))
    if element_type == "text":
        text = element.get("text", "")
        if isinstance(text, list):
            font_size = fitted_font_size(text, width=width, height=height, preferred=int(typography.get("body_size", 15)), minimum=min_body)
            warn_if_dense(warnings, slide_index=index, label="text list", text=text, width=width, height=height, font_size=font_size, minimum=min_body)
            add_bullets(slide, [str(item) for item in text], left, top, width, height, style)
        else:
            font_size = fitted_font_size(str(text), width=width, height=height, preferred=int(typography.get("body_size", 15)), minimum=min_body)
            warn_if_dense(warnings, slide_index=index, label="text", text=str(text), width=width, height=height, font_size=font_size, minimum=min_body)
            add_textbox(slide, str(text), left, top, width, height, style, size=int(style.get("typography", {}).get("body_size", 15)))
    elif element_type == "bullets":
        items = [str(item) for item in element.get("items", [])]
        font_size = fitted_font_size(items, width=width, height=height, preferred=int(typography.get("body_size", 15)), minimum=min_body)
        warn_if_dense(warnings, slide_index=index, label="bullets", text=items, width=width, height=height, font_size=font_size, minimum=min_body)
        add_bullets(slide, items, left, top, width, height, style)
    elif element_type == "callout":
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(style, "surface")
        shape.line.color.rgb = color(style, "primary")
        callout_text = str(element.get("text", ""))
        inner_width = width - Inches(0.36)
        inner_height = height - Inches(0.24)
        font_size = fitted_font_size(callout_text, width=inner_width, height=inner_height, preferred=15, minimum=min_body)
        warn_if_dense(warnings, slide_index=index, label="callout", text=callout_text, width=inner_width, height=inner_height, font_size=font_size, minimum=min_body)
        add_textbox(slide, callout_text, left + Inches(0.18), top + Inches(0.12), inner_width, inner_height, style, size=15, bold=True)
    elif element_type == "table":
        add_table(slide, element.get("rows", []), left, top, width, height, style)
    elif element_type == "chart_placeholder":
        add_placeholder(slide, str(element.get("title", "Chart")), left, top, width, height, style, kind=str(element.get("chart_type", "Chart")))
    elif element_type in {"image", "picture", "image_placeholder"}:
        path_value = image_path_from_element(element)
        resolved = resolve_local_path(path_value, base_dir=base_dir)
        if resolved:
            add_image(slide, resolved, left, top, width, height, style, caption=str(element.get("caption") or ""))
        else:
            label = str(element.get("label") or element.get("alt") or path_value or "Image")
            add_placeholder(slide, label, left, top, width, height, style, kind="Image")
            warnings.append(f"Slide {index}: image not found; rendered placeholder for `{label}`")
    elif element_type in {"video", "video_link", "video_placeholder"}:
        add_video_link(slide, element, left, top, width, height, style, base_dir=base_dir)
        if not (element.get("video_url") or element.get("url") or element.get("href") or element.get("video_path")):
            warnings.append(f"Slide {index}: video element has no video_url/video_path")
    else:
        warnings.append(f"Slide {index}: unsupported element type `{element_type}`")


def render_slide(
    prs: Presentation,
    slide_data: dict[str, Any],
    style: dict[str, Any],
    index: int,
    total: int,
    warnings: list[str],
    *,
    base_dir: Path,
) -> None:
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    layout = str(slide_data.get("layout", "content"))
    headline = str(slide_data.get("headline") or slide_data.get("title") or f"Slide {index}")
    preserve_background = bool(template_fidelity(style).get("preserve_background")) and use_template_chrome(style)
    if not preserve_background:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color(style, "background")

    typography = style.get("typography", {})
    tokens = style.get("layout_tokens", {})
    margin_left = float(tokens.get("margin_left", 0.65))
    margin_right = float(tokens.get("margin_right", 0.65))
    gutter = float(tokens.get("gutter", 0.28))
    header_h = float(tokens.get("header_height", 0.25))
    footer_h = float(tokens.get("footer_height", 0.28))

    if layout in {"title", "cover"}:
        add_design_chrome(slide, style, index, total, headline)
        title_box = box_for_role(style, layout, "headline") or box_for_role(style, layout, "title")
        title_top = 1.4 + max(0, header_h - 0.25)
        if title_box:
            add_textbox(slide, headline, *title_box, style, role="title", size=int(typography.get("title_size", 34)))
        else:
            add_textbox(slide, headline, inches(0.85), inches(title_top), inches(11.5), inches(1.05), style, role="title", size=int(typography.get("title_size", 34)))
        subtitle = slide_data.get("subtitle") or slide_data.get("purpose")
        if subtitle:
            subtitle_box = box_for_role(style, layout, "subtitle")
            if subtitle_box:
                add_textbox(slide, str(subtitle), *subtitle_box, style, size=int(typography.get("subtitle_size", 18)))
            else:
                add_textbox(slide, str(subtitle), inches(0.9), inches(title_top + 1.08), inches(10.9), inches(0.72), style, size=int(typography.get("subtitle_size", 18)))
    elif layout in {"section", "divider"}:
        add_design_chrome(slide, style, index, total, headline)
        title_box = box_for_role(style, layout, "headline") or box_for_role(style, layout, "title")
        if title_box:
            add_textbox(slide, headline, *title_box, style, role="title", size=min(int(typography.get("title_size", 32)), 34))
        else:
            add_textbox(slide, headline, inches(0.85), inches(2.25), inches(11.7), inches(1.0), style, role="title", size=min(int(typography.get("title_size", 32)), 34))
        section = slide_data.get("section")
        if section:
            add_textbox(slide, str(section), inches(0.9), inches(1.72), inches(5.6), inches(0.32), style, role="caption", size=11)
    else:
        add_design_chrome(slide, style, index, total, headline)
        title_top = max(0.34, header_h + 0.12)
        title_box = box_for_role(style, layout, "headline") or box_for_role(style, layout, "title")
        if title_box:
            add_textbox(slide, headline, *title_box, style, role="title", size=min(int(typography.get("title_size", 28)), 28))
        else:
            add_textbox(slide, headline, inches(margin_left), inches(title_top), inches(13.333 - margin_left - margin_right), inches(0.58), style, role="title", size=min(int(typography.get("title_size", 28)), 28))
        elements = [item for item in slide_data.get("elements", []) if isinstance(item, dict)]
        boxes = element_boxes_from_template(style, layout, elements) or content_boxes(
            len(elements),
            top=title_top + 0.82,
            bottom=7.5 - footer_h - 0.34,
            left=margin_left + 0.1,
            right=13.333 - margin_right - 0.1,
            gutter=gutter,
        )
        for element, box in zip(elements, boxes):
            if not isinstance(element, dict):
                continue
            render_element(slide, element, box, style, warnings, index, base_dir=base_dir)

    source_notes = slide_data.get("source_notes", [])
    if source_notes:
        source_text = "; ".join(str(item) for item in source_notes)
        add_textbox(slide, f"Source: {source_text}", inches(margin_left), inches(6.72), inches(10.8), inches(0.2), style, role="footer", size=7)

    add_notes(slide, str(slide_data.get("speaker_notes", "")))


def render(deck_spec: dict[str, Any], style_spec: dict[str, Any], output: Path, report: Path, *, base_dir: Path) -> None:
    style = merge_style(style_spec)
    warnings: list[str] = []
    template_path = resolve_style_path(str(style.get("template_pptx") or ""), base_dir=base_dir)
    if template_path:
        prs = Presentation(str(template_path))
        slide_id_list = prs.slides._sldIdLst  # python-pptx has no public API for clearing template slides.
        for slide_id in list(slide_id_list):
            slide_id_list.remove(slide_id)
    else:
        prs = Presentation()
        if style.get("template_pptx"):
            warnings.append(f"Template PPTX not found: {style.get('template_pptx')}")

    if style.get("slide_size") == "16:9":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slides = deck_spec.get("slides", [])
    if not isinstance(slides, list) or not slides:
        raise ValueError("deck_spec must contain a non-empty slides list")

    for index, slide_data in enumerate(slides, start=1):
        if not isinstance(slide_data, dict):
            warnings.append(f"Slide {index}: skipped non-object slide")
            continue
        render_slide(prs, slide_data, style, index, len(slides), warnings, base_dir=base_dir)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        "\n".join(
            [
                "# PPTX Encoding Report",
                "",
                f"Output: `{output}`",
                f"Slides: {len(slides)}",
                "Editable strategy: native PowerPoint text boxes, shapes, tables and notes.",
                "",
                "## Warnings",
                *(f"- {warning}" for warning in warnings),
                "" if warnings else "- None",
                "",
            ]
        ),
        encoding="utf-8",
    )


def main() -> int:
    parser = argparse.ArgumentParser(description="Render editable PPTX from deck and style specs.")
    parser.add_argument("--deck-spec", required=True, type=Path)
    parser.add_argument("--style-spec", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--report", required=True, type=Path)
    args = parser.parse_args()

    try:
        deck_spec = load_json(args.deck_spec)
        style_spec = load_json(args.style_spec)
        render(deck_spec, style_spec, args.output, args.report, base_dir=args.deck_spec.parent)
    except Exception as exc:
        print(f"render_deck.py: error: {exc}", file=sys.stderr)
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
