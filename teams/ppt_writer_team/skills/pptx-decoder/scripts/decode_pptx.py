#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
from collections import Counter
from pathlib import Path
from statistics import median
from typing import Any


def fail_missing_dependency() -> None:
    print(
        "Missing dependency: python-pptx. Install with `pip install python-pptx`.",
        file=sys.stderr,
    )
    raise SystemExit(2)


try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    fail_missing_dependency()


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    try:
        return str(rgb)
    except Exception:
        return None


def emu_to_inches(value: int) -> float:
    return round(value / 914400, 3)


def safe_color(color_format) -> str | None:
    try:
        if not color_format or not color_format.type:
            return None
        rgb = getattr(color_format, "rgb", None)
        if rgb:
            return rgb_to_hex(rgb)
        theme_color = getattr(color_format, "theme_color", None)
        if theme_color:
            return f"theme:{theme_color}"
    except Exception:
        return None
    return None


def region_for_shape(shape, slide_width: int, slide_height: int) -> str:
    center_y = shape.top + shape.height / 2
    center_x = shape.left + shape.width / 2
    if center_y < slide_height * 0.14:
        return "header"
    if center_y > slide_height * 0.86:
        return "footer"
    if center_x < slide_width * 0.38:
        return "left_body"
    if center_x > slide_width * 0.62:
        return "right_body"
    return "center_body"


def text_runs(shape) -> list[dict[str, object]]:
    if not getattr(shape, "has_text_frame", False):
        return []

    runs: list[dict[str, object]] = []
    for paragraph in shape.text_frame.paragraphs:
        paragraph_text = paragraph.text.strip()
        if not paragraph_text:
            continue
        for run in paragraph.runs:
            font = run.font
            runs.append(
                {
                    "text": run.text,
                    "font_name": font.name,
                    "font_size_pt": font.size.pt if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                    "color": safe_color(font.color),
                }
            )
        if not paragraph.runs:
            runs.append({"text": paragraph_text})
    return runs


def shape_info(shape) -> dict[str, object]:
    info: dict[str, object] = {
        "shape_type": str(shape.shape_type),
        "name": shape.name,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }

    if getattr(shape, "has_text_frame", False):
        info["text"] = shape.text.strip()
        info["runs"] = text_runs(shape)
        info["text_margin_left"] = shape.text_frame.margin_left
        info["text_margin_right"] = shape.text_frame.margin_right
        info["text_margin_top"] = shape.text_frame.margin_top
        info["text_margin_bottom"] = shape.text_frame.margin_bottom

    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        info["table"] = rows

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["picture"] = True

    if getattr(shape, "has_chart", False):
        info["chart"] = True
        info["chart_type"] = str(shape.chart.chart_type)

    try:
        fill = shape.fill
        if fill and fill.type:
            info["fill_type"] = str(fill.type)
            info["fill_color"] = safe_color(fill.fore_color)
    except Exception:
        pass

    try:
        line = shape.line
        if line and line.color and line.color.type:
            info["line_color"] = safe_color(line.color)
            info["line_width"] = line.width
    except Exception:
        pass

    return info


def slide_background_color(slide) -> str | None:
    try:
        fill = slide.background.fill
        return safe_color(fill.fore_color)
    except Exception:
        return None


def notes_text(slide) -> str:
    try:
        notes_slide = slide.notes_slide
        return notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def top_hex(counter: Counter[str], *, exclude: set[str] | None = None, limit: int = 8) -> list[str]:
    exclude = exclude or set()
    values = []
    for value, _count in counter.most_common():
        if value.startswith("theme:"):
            continue
        normalized = value.upper().lstrip("#")
        if len(normalized) == 6 and normalized not in exclude and normalized not in values:
            values.append(normalized)
        if len(values) >= limit:
            break
    return values


def hex_luminance(value: str) -> float:
    value = value.upper().lstrip("#")
    if len(value) != 6:
        return 1.0
    r = int(value[0:2], 16) / 255
    g = int(value[2:4], 16) / 255
    b = int(value[4:6], 16) / 255
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def role_for_text_box(item: dict[str, Any], slide_width: int, slide_height: int) -> str:
    text = str(item.get("text") or "").strip()
    height = int(item.get("height") or 0)
    width = int(item.get("width") or 0)
    top = int(item.get("top") or 0)
    if top < slide_height * 0.16 and len(text) <= 90 and height < slide_height * 0.22:
        return "header"
    if top > slide_height * 0.82:
        return "footer"
    if len(text) <= 80 and height < slide_height * 0.25 and width > slide_width * 0.35:
        return "title"
    if len(text) <= 120 and height < slide_height * 0.18:
        return "caption"
    return "body"


def style_spec_suggestion(
    *,
    slide_width: int,
    slide_height: int,
    font_counter: Counter[str],
    color_counter: Counter[str],
    font_color_counter: Counter[str],
    size_counter: Counter[int],
    region_counter: Counter[str],
    background_counter: Counter[str],
    fill_counter: Counter[str],
    footer_counter: Counter[str],
    header_counter: Counter[str],
) -> dict[str, Any]:
    background_candidates = top_hex(background_counter, limit=4)
    fill_colors = top_hex(fill_counter, exclude=set(background_candidates), limit=8)
    all_colors = top_hex(color_counter, exclude={"FFFFFF", "000000"}, limit=8)
    colors = fill_colors + [item for item in all_colors if item not in fill_colors]
    background = "FFFFFF"
    for candidate in background_candidates:
        candidate = candidate.upper().lstrip("#")
        if hex_luminance(candidate) > 0.72:
            background = candidate
            break
    else:
        if background_candidates:
            background = background_candidates[0]

    text = "1F2937"
    text_candidates = top_hex(font_color_counter, limit=8)
    dark_text_candidates = [candidate for candidate in text_candidates if hex_luminance(candidate) < 0.45]
    if dark_text_candidates:
        text = min(dark_text_candidates, key=hex_luminance)
    elif "000000" in font_color_counter:
        text = "000000"
    elif text_candidates:
        text = text_candidates[0]

    accent_colors = [candidate for candidate in colors if candidate not in {background, text}]
    muted = "6B7280"
    if text_candidates:
        muted = max(text_candidates, key=hex_luminance)
    surface_candidates = [candidate for candidate in accent_colors if hex_luminance(candidate) > 0.78]
    surface = surface_candidates[0] if surface_candidates else "F3F4F6"

    font = font_counter.most_common(1)[0][0] if font_counter else "Arial"
    sizes = [size for size, _count in size_counter.most_common()]
    title_size = max(sizes) if sizes else 30
    body_candidates = [size for size in sizes if 10 <= size <= min(title_size, 24)]
    body_size = int(median(body_candidates[:5])) if body_candidates else 16

    return {
        "slide_size": "16:9" if abs((slide_width / slide_height) - (16 / 9)) < 0.08 else "custom",
        "slide_width_emu": slide_width,
        "slide_height_emu": slide_height,
        "theme_colors": {
            "background": background,
            "text": text,
            "primary": accent_colors[0] if accent_colors else "2563EB",
            "secondary": accent_colors[1] if len(accent_colors) > 1 else "14B8A6",
            "accent": accent_colors[2] if len(accent_colors) > 2 else "F59E0B",
            "muted": muted,
            "surface": surface,
        },
        "fonts": {"title": font, "body": font, "caption": font},
        "typography": {
            "title_size": max(24, min(int(title_size), 40)),
            "subtitle_size": max(16, min(body_size + 3, 24)),
            "body_size": max(12, min(body_size, 20)),
            "caption_size": max(8, min(body_size - 3, 12)),
            "min_body_size": 10,
            "min_title_size": 20,
        },
        "layout_tokens": {
            "margin_left": 0.65,
            "margin_right": 0.65,
            "header_height": 0.42 if region_counter.get("header") else 0.25,
            "footer_height": 0.34 if region_counter.get("footer") else 0.28,
            "gutter": 0.28,
        },
        "template_patterns": {
            "regions": region_counter.most_common(),
            "header_text_examples": [text for text, _ in header_counter.most_common(3)],
            "footer_text_examples": [text for text, _ in footer_counter.most_common(3)],
            "uses_header": bool(region_counter.get("header")),
            "uses_footer": bool(region_counter.get("footer")),
            "style_transfer_note": "Transfer colors, typography, spacing, headers, footers and visual rhythm only; do not reuse template business content.",
        },
    }


def decode(path: Path, mode: str) -> dict[str, object]:
    prs = Presentation(path)
    font_counter: Counter[str] = Counter()
    color_counter: Counter[str] = Counter()
    font_color_counter: Counter[str] = Counter()
    size_counter: Counter[int] = Counter()
    layout_counter: Counter[str] = Counter()
    region_counter: Counter[str] = Counter()
    fill_counter: Counter[str] = Counter()
    line_counter: Counter[str] = Counter()
    background_counter: Counter[str] = Counter()
    header_counter: Counter[str] = Counter()
    footer_counter: Counter[str] = Counter()
    slides = []
    warnings: list[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name if slide.slide_layout else ""
        layout_counter[layout_name] += 1

        shapes = [shape_info(shape) for shape in slide.shapes]
        texts = []
        for item in shapes:
            item["region"] = region_for_shape(type("ShapeLike", (), item), prs.slide_width, prs.slide_height)
            region_counter[str(item["region"])] += 1
            text = item.get("text")
            if isinstance(text, str) and text:
                texts.append(text)
                role = role_for_text_box(item, prs.slide_width, prs.slide_height)
                item["text_role_guess"] = role
                if role == "header":
                    header_counter[text[:120]] += 1
                elif role == "footer":
                    footer_counter[text[:120]] += 1
            for run in item.get("runs", []) if isinstance(item.get("runs"), list) else []:
                font = run.get("font_name")
                color = run.get("color")
                size = run.get("font_size_pt")
                if font:
                    font_counter[str(font)] += 1
                if color:
                    color_counter[str(color)] += 1
                    font_color_counter[str(color)] += 1
                if isinstance(size, (int, float)):
                    size_counter[int(round(size))] += 1
            fill_color = item.get("fill_color")
            line_color = item.get("line_color")
            if fill_color:
                color_counter[str(fill_color)] += 1
                fill_counter[str(fill_color)] += 1
            if line_color:
                color_counter[str(line_color)] += 1
                line_counter[str(line_color)] += 1

        bg_color = slide_background_color(slide)
        if bg_color:
            color_counter[bg_color] += 2
            background_counter[bg_color] += 1

        slide_record: dict[str, object] = {
            "slide_number": index,
            "layout": layout_name,
            "shape_count": len(shapes),
            "text_blocks": texts,
            "notes": notes_text(slide),
            "background_color": bg_color,
        }
        if mode == "template":
            slide_record["shapes"] = [
                {
                    key: value
                    for key, value in item.items()
                    if key
                    in {
                        "shape_type",
                        "name",
                        "left",
                        "top",
                        "width",
                        "height",
                        "fill_color",
                        "fill_type",
                        "line_color",
                        "line_width",
                        "runs",
                        "region",
                        "text_role_guess",
                    }
                }
                for item in shapes
            ]
        else:
            slide_record["shapes"] = shapes
        slides.append(slide_record)

    if mode == "template":
        warnings.append("Template mode extracts style signals only. Do not reuse unrelated template content.")

    return {
        "file": str(path),
        "mode": mode,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_count": len(prs.slides),
        "slides": slides,
        "style_summary": {
            "layouts": layout_counter.most_common(),
            "fonts": font_counter.most_common(12),
            "font_sizes_pt": size_counter.most_common(12),
            "colors": color_counter.most_common(16),
            "fill_colors": fill_counter.most_common(12),
            "line_colors": line_counter.most_common(12),
            "font_colors": font_color_counter.most_common(12),
            "background_colors": background_counter.most_common(8),
            "regions": region_counter.most_common(),
            "headers": header_counter.most_common(8),
            "footers": footer_counter.most_common(8),
            "slide_size_inches": {
                "width": emu_to_inches(prs.slide_width),
                "height": emu_to_inches(prs.slide_height),
            },
        },
        "style_spec_suggestion": style_spec_suggestion(
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            font_counter=font_counter,
            color_counter=color_counter,
            font_color_counter=font_color_counter,
            size_counter=size_counter,
            region_counter=region_counter,
            background_counter=background_counter,
            fill_counter=fill_counter,
            footer_counter=footer_counter,
            header_counter=header_counter,
        )
        if mode == "template"
        else None,
        "warnings": warnings,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Decode PPTX content or template style into JSON.")
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--mode", choices=("template", "material"), default="template")
    args = parser.parse_args()

    if not args.input.exists():
        print(f"Input not found: {args.input}", file=sys.stderr)
        return 2

    result = decode(args.input, args.mode)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
